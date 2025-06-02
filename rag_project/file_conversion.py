import os
import PyPDF2
import pandas as pd
from docx import Document
import pptx
import csv

def convert_to_text(file_path, file_ext):
    """Convert various file formats to plain text"""
    text = ""
    
    try:
        if file_ext == '.pdf':
            with open(file_path, 'rb') as f:
                reader = PyPDF2.PdfReader(f)
                for page in reader.pages:
                    text += page.extract_text() + "\n"
                    
        elif file_ext in ('.doc', '.docx'):
            doc = Document(file_path)
            for para in doc.paragraphs:
                text += para.text + "\n"
                
        elif file_ext in ('.xls', '.xlsx'):
            df = pd.read_excel(file_path)
            text = df.to_string()
            
        elif file_ext == '.csv':
            df = pd.read_csv(file_path)
            text = df.to_string()
            
        elif file_ext in ('.ppt', '.pptx'):
            prs = pptx.Presentation(file_path)
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"
                        
        elif file_ext == '.txt':
            with open(file_path, 'r', encoding='utf-8') as f:
                text = f.read()
                
    except Exception as e:
        raise Exception(f"Error converting {file_ext} file: {str(e)}")
    
    return text.strip() 